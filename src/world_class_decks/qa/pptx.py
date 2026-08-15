from __future__ import annotations

from collections.abc import Iterator
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from world_class_decks.models import AuditReport, Finding, Severity
from world_class_decks.qa.copy import find_copy_issues

PLACEHOLDERS = ("lorem ipsum", "todo", "tbd", "insert text", "placeholder", "your text here")


def _iter_text(slide_number: int, slide) -> Iterator[tuple[int, str, str | None]]:
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            text = shape.text or ""
            if text.strip():
                yield slide_number, text, getattr(shape, "name", None)
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for child in shape.shapes:
                if getattr(child, "has_text_frame", False) and (child.text or "").strip():
                    yield slide_number, child.text, getattr(child, "name", None)


def _overlap_ratio(a, b) -> float:
    x1 = max(a.left, b.left)
    y1 = max(a.top, b.top)
    x2 = min(a.left + a.width, b.left + b.width)
    y2 = min(a.top + a.height, b.top + b.height)
    if x2 <= x1 or y2 <= y1:
        return 0.0
    inter = (x2 - x1) * (y2 - y1)
    denom = min(a.width * a.height, b.width * b.height)
    return float(inter / denom) if denom else 0.0


def audit_pptx(path: Path, *, max_words_per_slide: int = 90, min_font_pt: float = 14.0) -> AuditReport:
    prs = Presentation(str(path))
    findings: list[Finding] = []
    texts: list[tuple[int, str, str | None]] = []

    for slide_number, slide in enumerate(prs.slides, start=1):
        slide_words = 0
        shapes = list(slide.shapes)
        for shape in shapes:
            name = getattr(shape, "name", None)
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > prs.slide_width or shape.top + shape.height > prs.slide_height:
                findings.append(Finding("off-slide", "Shape extends outside slide bounds.", Severity.ERROR, slide_number, name))

            if getattr(shape, "has_text_frame", False):
                text = shape.text or ""
                slide_words += len(text.split())
                lowered = text.casefold()
                for marker in PLACEHOLDERS:
                    if marker in lowered:
                        findings.append(Finding("placeholder", f"Placeholder text contains {marker!r}.", Severity.ERROR, slide_number, name))
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        size = run.font.size
                        if size is not None and size.pt < min_font_pt and text.strip():
                            findings.append(
                                Finding(
                                    "tiny-text",
                                    f"Text run is {size.pt:.1f}pt (< {min_font_pt:.1f}pt).",
                                    Severity.WARNING,
                                    slide_number,
                                    name,
                                )
                            )

        if slide_words > max_words_per_slide:
            findings.append(
                Finding(
                    "dense-slide",
                    f"Slide has {slide_words} words (> {max_words_per_slide}).",
                    Severity.WARNING,
                    slide_number,
                )
            )

        # Conservative geometry heuristic: only flag strong overlaps involving text.
        for i, a in enumerate(shapes):
            for b in shapes[i + 1 :]:
                if not (getattr(a, "has_text_frame", False) or getattr(b, "has_text_frame", False)):
                    continue
                ratio = _overlap_ratio(a, b)
                if ratio >= 0.60:
                    findings.append(
                        Finding(
                            "strong-overlap",
                            f"Shapes overlap substantially ({ratio:.0%}). Verify rendered slide.",
                            Severity.WARNING,
                            slide_number,
                            f"{getattr(a, 'name', '?')} ↔ {getattr(b, 'name', '?')}",
                        )
                    )

        texts.extend(_iter_text(slide_number, slide))

    findings.extend(find_copy_issues(texts))
    return AuditReport(file=path, slide_count=len(prs.slides), findings=findings)
