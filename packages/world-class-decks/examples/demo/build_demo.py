from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

OUT = Path(__file__).with_name("demo.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BG = RGBColor(245, 243, 238)
INK = RGBColor(24, 24, 24)
MUTED = RGBColor(90, 90, 90)
ACCENT = RGBColor(183, 69, 48)


def bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def text(slide, value, x, y, w, h, size, bold=False, color=INK, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = value
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


# 1 — thesis
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)
text(slide, "Fewer incidents. More impact.", 0.8, 1.3, 11.7, 0.8, 34, True)
text(slide, "The reliability problem moved from frequency to containment speed.", 0.82, 2.25, 8.4, 0.6, 18, False, MUTED)
accent = slide.shapes.add_shape(1, Inches(0.82), Inches(3.4), Inches(2.25), Inches(0.09))
accent.fill.solid()
accent.fill.fore_color.rgb = ACCENT
accent.line.fill.background()

# 2 — comparison
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)
text(slide, "Incident count fell 42% — but two long events erased the gain", 0.8, 0.55, 11.8, 0.7, 26, True)
text(slide, "19", 1.2, 2.0, 2.0, 0.8, 48, True, ACCENT)
text(slide, "Sev-1 incidents\nprior quarter", 1.2, 2.9, 2.4, 0.9, 17, False, MUTED)
text(slide, "11", 5.1, 2.0, 2.0, 0.8, 48, True, INK)
text(slide, "Sev-1 incidents\ncurrent quarter", 5.1, 2.9, 2.4, 0.9, 17, False, MUTED)
text(slide, "+14%", 9.0, 2.0, 2.2, 0.8, 48, True, ACCENT)
text(slide, "customer-impact\nminutes", 9.0, 2.9, 2.5, 0.9, 17, False, MUTED)

# 3 — decision
slide = prs.slides.add_slide(prs.slide_layouts[6]); bg(slide)
text(slide, "Fund containment speed, not another incident-count program", 0.8, 0.7, 11.6, 0.7, 28, True)
text(slide, "Decision", 0.82, 2.0, 2.0, 0.4, 16, True, ACCENT)
text(slide, "Prioritize automated diagnosis, blast-radius controls, and rehearsed rollback paths.", 0.82, 2.55, 8.4, 1.2, 25, False)
text(slide, "Success measure", 0.82, 4.55, 2.4, 0.4, 16, True, ACCENT)
text(slide, "Reduce median customer-impact minutes per Sev-1 by 35% next quarter.", 0.82, 5.1, 9.3, 0.8, 22, False)

prs.save(OUT)
print(OUT)
