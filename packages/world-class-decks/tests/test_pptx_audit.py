from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from world_class_decks.qa.pptx import audit_pptx


def make_pptx(path: Path, text: str = "Specific claim") -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(24)
    prs.save(path)


def test_clean_pptx_passes(tmp_path):
    path = tmp_path / "clean.pptx"
    make_pptx(path)
    report = audit_pptx(path)
    assert report.passed


def test_placeholder_fails(tmp_path):
    path = tmp_path / "bad.pptx"
    make_pptx(path, "TODO replace this")
    report = audit_pptx(path)
    assert not report.passed
    assert any(f.code == "placeholder" for f in report.findings)


def test_large_empty_panel_fails(tmp_path):
    path = tmp_path / "empty-panel.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5), Inches(1), Inches(6), Inches(5))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(20, 60, 90)
    prs.save(path)
    report = audit_pptx(path)
    assert not report.passed
    assert any(f.code == "empty-panel" for f in report.findings)


def test_full_slide_background_is_not_an_empty_panel(tmp_path):
    path = tmp_path / "background.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(20, 60, 90)
    prs.save(path)

    report = audit_pptx(path)

    assert report.passed
    assert not any(f.code == "empty-panel" for f in report.findings)
