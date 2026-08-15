from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from world_class_decks.qa.capabilities import inspect_capabilities


def test_capability_inspection(tmp_path: Path):
    path = tmp_path / "deck.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "Hello"
    prs.save(path)
    report = inspect_capabilities(path)
    assert report["native"]["themes"] >= 1
    assert report["native"]["slide_masters"] >= 1
