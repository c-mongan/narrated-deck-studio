from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from world_class_decks.qa.pptx import audit_pptx


def make_pptx(path: Path, text: str = "Specific claim") -> None:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    run = box.text_frame.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(24)
    prs.save(path)


def test_clean_pptx_passes(tmp_path):
    path = tmp_path / "clean.pptx"
    make_pptx(path)
    report = audit_pptx(path)
    assert report.passed


def test_placeholder_fails(tmp_path):
    path = tmp_path / "bad.pptx"
    make_pptx(path, "TODO replace this")
    report = audit_pptx(path)
    assert not report.passed
    assert any(f.code == "placeholder" for f in report.findings)
